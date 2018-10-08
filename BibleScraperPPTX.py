# Ferdinand Mudjialim 10/08/2018
# Parses Bible verses using Beautiful Soup and requests module
# Scrapes the site below to get Terjemahan Baru Indonesian Bible verses
# Compiles everything into a PowerPoint with each entry as a separate slide
# http://www.bibledbdata.org/onlinebibles/indonesian_tb/01_001.htm

import requests, bs4, time
from pptx import Presentation
# from pptx.enum.text import MSO_AUTO_SIZE
# from pptx.util import Inches, Pt


def printVerses(bookName, chapterNumber, verseNumbers):

    bookList = ['kejadian', 'keluaran', 'imamat', 'bilangan', 'ulangan',
                'yosua', 'hakim-hakim', 'rut', '1samuel', '2samuel',
                '1raja-raja', '2raja-raja', '1tawarikh', '2tawarikh',
                'ezra', 'nehemia', 'ester', 'ayub', 'mazmur', 'amsal',
                'pengkhotbah', 'kidungagung', 'yesaya', 'yeremia', 'ratapan',
                'yehezkiel', 'daniel', 'hosea', 'yoel', 'amos', 'obaja',
                'yunus', 'mikha', 'nahum', 'habakuk', 'zefanya', 'hagai',
                'zakharia', 'maleakhi', 'matius', 'markus', 'lukas', 'yohanes',
                'kisahpararasul', 'roma', '1korintus', '2korintus',
                'galatia', 'efesus', 'filipi', 'kolose', '1tesalonika',
                '2tesalonika', '1timotius', '2timotius', 'titus', 'filemon',
                'ibrani', 'yakobus', '1petrus', '2petrus', '1yohanes',
                '2yohanes', '3yohanes', 'yudas', 'wahyu']
    bookLookup = {'2samuel': '2 samuel', '1raja-raja': '1 raja-raja',
                  '2raja-raja': '2 raja-raja', '1tawarikh': '1 tawarikh',
                  '2tawarikh': '2 tawarikh', 'kidungagung': 'kidung agung',
                  'kisahpararasul': 'kisah para rasul',
                  '1korintus': '1 korintus', '2korintus': '2 korintus',
                  '1tesalonika': '1 tesalonika', '2tesalonika': '2 tesalonika',
                  '1timotius': '1 timotius', '2timotius': '2 timotius',
                  '1petrus': '1 petrus', '2petrus': '2 petrus',
                  '1yohanes': '1 yohanes', '2yohanes': '2 yohanes',
                  '3yohanes': '3 yohanes'}
    bookNumbers = ['0' + str(n) for n in range(1, 10)] \
                  + [str(m) for m in range(10, 67)]
    bookDict = dict(zip(bookList, bookNumbers))

    book = str(bookName.strip().lower())  # 2 digits format ex: 01
    chapter = str(chapterNumber.strip())  # 3 digits format ex: 001
    if len(chapter) <= 2:  # correct formatting for chapter numbers
        if len(chapter) == 1:
            chapter = '00' + chapter
        else:
            chapter = '0' + chapter
    verseInput = verseNumbers.strip()


    res = requests.get('http://www.bibledbdata.org/onlinebibles/indonesian_tb/'
                       + str(bookDict[book])
                       + '_'
                       + str(chapter)
                       + '.htm')
    res.raise_for_status()

    bibleSoup = bs4.BeautifulSoup(res.text, 'html.parser')
    bibleSoup = bibleSoup.select('blockquote')
    string = bibleSoup[0].getText()
    verseList = string.split('\n')
    verseList = [y for y in verseList if y != '']  # list of verses in order

    if verseInput == '*':  # prints all verseInput
        verseFrom = 0
        verseTo = len(verseList)
    else:
        parsedVerseNumbers = verseInput.split('-')
        # if verseInput in form 'int-'
        if verseInput[len(verseInput) - 1] == '-':
            verseFrom = int(parsedVerseNumbers[0]) - 1
            verseTo = len(verseList)
        # if verseInput in form '-int'
        elif verseInput[0] == '-':
            verseFrom = 0
            verseTo = int(parsedVerseNumbers[1])
        # if verseInput in form 'int' and '-' is not found
        elif parsedVerseNumbers[0].isdigit() and verseInput.find('-') == -1:
            verseFrom = int(parsedVerseNumbers[0]) - 1
            verseTo = verseFrom + 1
        # normal from-to format
        else:
            verseFrom = int(parsedVerseNumbers[0]) - 1
            verseTo = int(parsedVerseNumbers[1])

    outputString = ''
    if book in bookLookup.keys():  # if the book title needs to be 'decoded'
        outputString += bookLookup[book].title() + ' ' \
                        + chapter.lstrip('0')\
                        + '\n'
    else:
        outputString += book.title() + ' ' \
                        + chapter.lstrip('0') \
                        + '\n'
    for i in range(verseFrom, verseTo):
        outputString += verseList[i] + '\n'
    return outputString


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = 'IBC'
subtitle.text = time.strftime('%m/%d/%Y')  # auto date on subtitle


with open('input.txt') as fp:
        line = fp.readline()
        while line != '':
            lineParseList = line.split()

            if lineParseList[0] == '#':
                line = fp.readline()
            else:
                # add slide with text and title.
                content = printVerses(lineParseList[0], lineParseList[1],
                                      lineParseList[2])

                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes

                title_shape = shapes.title
                body_shape = shapes.placeholders[1]

                title_shape.text = content.split('\n')[0]  # print title

                tf = body_shape.text_frame
                allVerses = '\n'.join(content.split('\n')[1:])
                allVersesList = allVerses.split('\n')

                tf.text = allVersesList[0]  # add first verse bullet

                for i in range(1, len(allVersesList)):  # add verse bullets
                    p = tf.add_paragraph()
                    p.text = allVersesList[i]

                tf.fit_text(font_family='Calibri', max_size=32, bold=False,
                            italic=False, font_file=None)
                # tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                line = fp.readline()

prs.save(time.strftime('%Y-%m-%d IBC.pptx'))  # saves file with date as name
